from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colours
CREAM = RGBColor(0xFA, 0xF3, 0xE0)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xD4, 0xED, 0xE9)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name='Calibri',
                font_size=14, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox, tf


def add_heading(slide, top, text, font_size=36):
    txBox, tf = add_textbox(slide, 0.8, top, 11.7, 0.8, text,
                            'Arial', font_size, True, NAVY, PP_ALIGN.LEFT)
    return txBox, tf


def add_accent_bar(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================
# SLIDE 1: Title
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)
add_accent_bar(slide, 0, 7.35, 13.333, 0.15, CORAL)

add_textbox(slide, 1.5, 2.0, 10.3, 1.2,
            "Using PiXL Resources\nin Your Lesson",
            'Arial', 42, True, NAVY, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.8, 10.3, 0.6,
            "HAB Crayford Science Department",
            'Calibri', 22, False, TEAL, PP_ALIGN.CENTER)

# Small decorative line
add_accent_bar(slide, 5.5, 3.5, 2.3, 0.04, CORAL)

add_notes(slide, "This is a quick guide to help you use PiXL resources effectively. "
                 "Don't worry about learning everything at once - start with one topic and one lesson.")


# ============================
# SLIDE 2: The 6 Resources You Need
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "The 6 Resources You Need")
add_accent_bar(slide, 0.8, 1.1, 3.0, 0.04, CORAL)

resources = [
    ("KnowIT PPT", "Teacher presentation (key knowledge)"),
    ("KnowIT Questions", "Student recall quiz"),
    ("GraspIT Answers", "Mark scheme (project on board)"),
    ("GraspIT Write On", "Student worksheet (print)"),
    ("GraspIT Reusable", "Laminate version"),
    ("ThinkIT", "Challenge questions (grades 7-9)"),
]

y_start = 1.6
for i, (name, desc) in enumerate(resources):
    y = y_start + i * 0.85
    # Teal rounded rect for resource name
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(y), Inches(3.2), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Description
    add_textbox(slide, 4.5, y + 0.05, 7.5, 0.55, "= " + desc, 'Calibri', 16, False, NAVY)

add_notes(slide, "You don't need all 6 every lesson. Start with KnowIT Questions for your "
                 "Do Now and GraspIT Write On for independent practice. Add more as you get comfortable.")


# ============================
# SLIDE 3: Your Lesson Structure
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "Your Lesson Structure")
add_accent_bar(slide, 0.8, 1.1, 3.0, 0.04, CORAL)

phases = [
    ("Do Now", "5 min", TEAL),
    ("Vocab", "3 min", CORAL),
    ("I Do", "10 min", NAVY),
    ("We Do", "10 min", NAVY),
    ("You Do", "10 min", TEAL),
    ("Exit", "5 min", CORAL),
]

box_w = 1.6
gap = 0.35
total_w = len(phases) * box_w + (len(phases) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (phase, time, color) in enumerate(phases):
    x = start_x + i * (box_w + gap)
    y = 2.8

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(box_w), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = phase
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = time
    p2.font.name = 'Calibri'
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # Arrow between boxes
    if i < len(phases) - 1:
        ax = x + box_w + 0.02
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(ax), Inches(y + 0.65), Inches(gap - 0.04), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        arrow.line.fill.background()

add_textbox(slide, 1.5, 5.2, 10.3, 0.6,
            "This follows the I Do, We Do, You Do model. "
            "You don't have to use every phase every lesson.",
            'Calibri', 14, False, NAVY, PP_ALIGN.CENTER)

add_notes(slide, "This follows the I Do, We Do, You Do model. The Do Now gets retrieval "
                 "going straight away. You don't have to use every phase every lesson.")


# ============================
# SLIDE 4: Do Now / Retrieval
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "Do Now / Retrieval (5 mins)")
add_accent_bar(slide, 0.8, 1.1, 3.5, 0.04, CORAL)

bullets = [
    "Project KnowIT Questions on the board",
    "Students answer from memory, no notes",
    "Mark together or swap and peer mark",
]

_, tf = add_textbox(slide, 1.2, 1.6, 10.5, 2.5, "", 'Calibri', 18)
tf.clear()
for j, b in enumerate(bullets):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = b
    run.font.name = 'Calibri'
    run.font.size = Pt(18)
    run.font.color.rgb = NAVY
    p.space_before = Pt(14)
    p.space_after = Pt(4)
    pPr = p._p.get_or_add_pPr()
    # Add bullet character
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    for existing in pPr.findall(qn('a:buNone')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)

# Tip box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.2), Inches(4.0), Inches(10.9), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_TEAL
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = ("Retrieval practice is the single most effective revision strategy. "
          "Even 5 minutes of recall at the start of a lesson makes a real difference.")
p.font.name = 'Calibri'
p.font.size = Pt(15)
p.font.color.rgb = NAVY
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

add_notes(slide, "This is the most important part. Retrieval practice is the single most "
                 "effective revision strategy. Even 5 minutes of recall at the start of a "
                 "lesson makes a real difference.")


# ============================
# SLIDE 5: I Do / We Do
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "I Do / We Do (20 mins)")
add_accent_bar(slide, 0.8, 1.1, 3.0, 0.04, CORAL)

# I Do box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.6), Inches(5.6), Inches(3.8))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.fill.background()

add_textbox(slide, 1.1, 1.8, 5.0, 0.5, "I Do (10 mins)", 'Arial', 24, True, TEAL)
_, tf = add_textbox(slide, 1.1, 2.5, 5.0, 2.5, "", 'Calibri', 16)
tf.clear()
ido_points = [
    "Pick 5-8 slides from the KnowIT PPT",
    "Talk through the key knowledge",
    "Don't show the whole thing!",
]
for j, b in enumerate(ido_points):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = b
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    p.space_before = Pt(10)
    pPr = p._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)

# We Do box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(1.6), Inches(5.6), Inches(3.8))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.fill.background()

add_textbox(slide, 7.2, 1.8, 5.0, 0.5, "We Do (10 mins)", 'Arial', 24, True, CORAL)
_, tf = add_textbox(slide, 7.2, 2.5, 5.0, 2.5, "", 'Calibri', 16)
tf.clear()
wedo_points = [
    "Project 2-3 GraspIT questions on the board",
    "Work through together as a class",
    "Use the Answers version to mark",
]
for j, b in enumerate(wedo_points):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = b
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    p.space_before = Pt(10)
    pPr = p._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)

# Warning box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xE8, 0xE1)  # light coral
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = ("The KnowIT PPTs are long. DO NOT try to go through the whole thing. "
          "Pick the slides that match what you're revising today.")
p.font.name = 'Calibri'
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = CORAL
p.alignment = PP_ALIGN.CENTER

add_notes(slide, "The KnowIT PPTs are long. DO NOT try to go through the whole thing. "
                 "Pick the slides that match what you're revising today. For We Do, pick "
                 "the questions students will find hardest.")


# ============================
# SLIDE 6: You Do / Exit
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "You Do / Exit (15 mins)")
add_accent_bar(slide, 0.8, 1.1, 3.0, 0.04, CORAL)

# You Do box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.6), Inches(5.6), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.fill.background()

add_textbox(slide, 1.1, 1.8, 5.0, 0.5, "You Do (10 mins)", 'Arial', 24, True, TEAL)
_, tf = add_textbox(slide, 1.1, 2.5, 5.0, 2.2, "", 'Calibri', 16)
tf.clear()
ydo_points = [
    "Hand out GraspIT Write On worksheets",
    "Students work independently",
    "They can self-mark using the Answers version",
]
for j, b in enumerate(ydo_points):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = b
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    p.space_before = Pt(10)
    pPr = p._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)

# Exit box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(1.6), Inches(5.6), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.fill.background()

add_textbox(slide, 7.2, 1.8, 5.0, 0.5, "Exit Ticket (5 mins)", 'Arial', 24, True, CORAL)
_, tf = add_textbox(slide, 7.2, 2.5, 5.0, 2.2, "", 'Calibri', 16)
tf.clear()
exit_points = [
    "Show 1-2 ThinkIT questions for stretch",
    "Or redo 3 KnowIT Questions to check recall stuck",
]
for j, b in enumerate(exit_points):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = b
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.color.rgb = NAVY
    p.space_before = Pt(10)
    pPr = p._p.get_or_add_pPr()
    for existing in pPr.findall(qn('a:buChar')):
        pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    pPr.append(buChar)

# Homework tip
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_TEAL
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = ("Homework tip: Print GraspIT Write On sheets for homework. "
          "Students complete at home, self-mark next lesson.")
p.font.name = 'Calibri'
p.font.size = Pt(15)
p.font.color.rgb = NAVY
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

add_notes(slide, "The Write On sheets are the best homework resource too. Print them, "
                 "students complete at home, self-mark next lesson. The ThinkIT is only "
                 "for your top students - don't stress about it if you're just starting out.")


# ============================
# SLIDE 7: How to Access Resources
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)

add_heading(slide, 0.4, "How to Access Resources")
add_accent_bar(slide, 0.8, 1.1, 3.5, 0.04, CORAL)

# URL box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2.5), Inches(1.6), Inches(8.3), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "missn77.github.io/hab-pixl-guide/"
p.font.name = 'Calibri'
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Password
add_textbox(slide, 2.5, 2.6, 8.3, 0.5,
            "Password:  habscience2026",
            'Calibri', 18, True, CORAL, PP_ALIGN.CENTER)

# Steps
steps = [
    ("1", "Pick your subject"),
    ("2", "Pick your topic"),
    ("3", "Preview or Download"),
    ("4", "Files download straight to your computer"),
]

for i, (num, text) in enumerate(steps):
    y = 3.5 + i * 0.8
    # Number circle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3.0), Inches(y), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 3.8, y + 0.05, 6.5, 0.45, text, 'Calibri', 18, False, NAVY)

add_textbox(slide, 2.5, 6.5, 8.3, 0.5,
            "No PiXL login needed. Bookmark the link.",
            'Calibri', 14, True, TEAL, PP_ALIGN.CENTER)

add_notes(slide, "No PiXL login needed. The app has every Edexcel resource organised by "
                 "topic. You can preview files in your browser before downloading. Bookmark the link.")


# ============================
# SLIDE 8: Start Small
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_accent_bar(slide, 0, 0, 13.333, 0.15, TEAL)
add_accent_bar(slide, 0, 7.35, 13.333, 0.15, CORAL)

add_heading(slide, 0.4, "Start Small")
add_accent_bar(slide, 0.8, 1.1, 1.8, 0.04, CORAL)

weeks = [
    ("This week", "Try a KnowIT Questions Do Now\nfor ONE topic", TEAL),
    ("Next week", "Add a GraspIT Write On\nfor independent practice", CORAL),
    ("Week after", "Try the full I Do, We Do,\nYou Do structure", NAVY),
]

for i, (week, desc, color) in enumerate(weeks):
    x = 1.0 + i * 3.9
    y = 1.8

    # Week label
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(3.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = week
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Description box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.7), Inches(3.5), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(24)

# Motivational footer
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_TEAL
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = ("Don't try to change everything at once. Start with just the Do Now.\n"
          "You'll be running full revision lessons within 3 weeks.")
p.font.name = 'Calibri'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(26)

add_notes(slide, "Don't try to change everything at once. Start with just the Do Now. "
                 "Once that feels natural, add the next piece. You'll be running full "
                 "revision lessons within 3 weeks.")


# ============================
# Clear metadata
# ============================
prs.core_properties.author = ''
prs.core_properties.comments = ''
prs.core_properties.title = 'Using PiXL Resources in Your Lesson'
prs.core_properties.subject = ''
prs.core_properties.keywords = ''
prs.core_properties.category = ''
prs.core_properties.last_modified_by = ''

output_path = '/Users/brendaneburagho/Desktop/PiXL_Teacher_Briefing.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
